from __future__ import annotations

import os
from pathlib import Path

from services.knowledge.ocr_vlm.ollama_vlm_adapter import OllamaVlmAdapter
from services.orchestrator.tools.base import ToolAdapter


class OcrProcessTool(ToolAdapter):
    id = "ocr.process"
    required_permissions = ["file.read"]
    declares_network_access = False

    def __init__(self):
        self._adapter = OllamaVlmAdapter(
            model=os.environ.get("VISION_MODEL_NAME")
            or os.environ.get("OCR_MODEL_NAME")
        )

    def invoke(self, inputs: dict) -> dict:
        path = inputs.get("path")
        if not path:
            raise ValueError("OcrProcessTool requires path")

        p = Path(path)
        suffix = p.suffix.lower()

        if suffix == ".pdf":
            return self._process_pdf(path, inputs)
        if suffix == ".pptx":
            return self._process_pptx(path)
        if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"}:
            return self._adapter.invoke(
                path,
                prompt=(
                    "Analyze this industrial image carefully. Perform OCR and "
                    "describe only observable text, labels, measurements, findings, "
                    "diagrams and other visible information. Do not invent facts."
                ),
            )

        raise ValueError(f"OCR/VLM does not support file type: {suffix or 'unknown'}")

    def _process_pdf(self, path: str, inputs: dict) -> dict:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError("PyMuPDF is required for PDF vision processing") from exc

        doc = fitz.open(path)
        contents: list[str] = []
        evidence: list[dict] = []
        requested_max = int(inputs.get("max_pages", 0) or 0)
        max_pages = len(doc) if requested_max <= 0 else min(len(doc), requested_max)

        for page_no in range(max_pages):
            page = doc[page_no]
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            image_path = Path(path).with_name(
                f"{Path(path).stem}__page_{page_no + 1}.png"
            )
            pix.save(str(image_path))
            try:
                result = self._adapter.invoke(
                    str(image_path),
                    prompt=(
                        "Read this scanned industrial document page. Extract all "
                        "visible text, labels, measurements, dates, findings and "
                        "other observable information. Do not invent missing content."
                    ),
                )
                text = str(result.get("content", "")).strip()
                if text:
                    contents.append(f"[Page {page_no + 1}]\n{text}")
                for item in result.get("evidence", []) or []:
                    item = dict(item)
                    item["page_or_region"] = f"page_{page_no + 1}"
                    item["source"] = path
                    evidence.append(item)
            finally:
                try:
                    image_path.unlink()
                except OSError:
                    pass

        return {
            "content": "\n\n".join(contents),
            "path": path,
            "model_id": self._adapter.id,
            "evidence": evidence,
        }

    def _process_pptx(self, path: str) -> dict:
        from pptx import Presentation

        prs = Presentation(path)
        blocks: list[str] = []
        evidence: list[dict] = []

        for slide_no, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text.strip()
                    if value:
                        texts.append(value)

                if getattr(shape, "shape_type", None) == 13:
                    image_path = Path(path).with_name(
                        f"{Path(path).stem}__slide_{slide_no}_{shape.shape_id}.png"
                    )
                    image_path.write_bytes(shape.image.blob)
                    try:
                        result = self._adapter.invoke(
                            str(image_path),
                            prompt=(
                                "Analyze this presentation image/diagram. Extract "
                                "visible text, labels, measurements, diagrams and "
                                "observable engineering findings. Do not invent facts."
                            ),
                        )
                        content = str(result.get("content", "")).strip()
                        if content:
                            texts.append(content)
                        for item in result.get("evidence", []) or []:
                            item = dict(item)
                            item["source"] = path
                            item["page_or_region"] = f"slide_{slide_no}"
                            evidence.append(item)
                    finally:
                        try:
                            image_path.unlink()
                        except OSError:
                            pass

            blocks.append(f"[Slide {slide_no}]\n" + "\n".join(texts))

        return {
            "content": "\n\n".join(blocks),
            "path": path,
            "model_id": self._adapter.id,
            "evidence": evidence,
        }
