from __future__ import annotations

import os
import tempfile
from pathlib import Path

from services.knowledge.rag.chunker import chunk_text
from services.knowledge.rag.store import MemoryVectorStore, VectorStore


class Retriever:
    def __init__(self, store: VectorStore):
        self._store = store

    def ingest_text(self, text: str, source: str, metadata: dict | None = None) -> int:
        chunks = chunk_text(text)
        self._store.add_chunks(chunks, source=source, metadata=metadata)
        return len(chunks)

    def ingest_file(self, path: str, metadata: dict | None = None) -> int:
        p = Path(path)
        suffix = p.suffix.lower()
        text = ""
        if suffix == ".pdf":
            import fitz
            doc = fitz.open(path)
            pages = [page.get_text("text") for page in doc]
            text = "\n\n".join(f"[Page {i+1}]\n{page}" for i, page in enumerate(pages))
            if len(text.strip()) < 40:
                from services.knowledge.ocr_vlm.ollama_vlm_adapter import OllamaVlmAdapter
                adapter = OllamaVlmAdapter()
                parts=[]
                for i, page in enumerate(doc):
                    pix=page.get_pixmap(matrix=fitz.Matrix(1.5,1.5), alpha=False)
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        tmp_path=Path(tmp.name)
                    try:
                        pix.save(str(tmp_path))
                        result=adapter.invoke(str(tmp_path), prompt="Extract all visible text, labels, numbers, findings and observable information from this page. Do not invent content.")
                        parts.append(f"[Page {i+1}]\n{result.get('content','')}")
                    finally:
                        tmp_path.unlink(missing_ok=True)
                text="\n\n".join(parts)
        elif suffix == ".docx":
            from docx import Document
            doc=Document(path)
            text="\n".join([p.text for p in doc.paragraphs] + [cell.text for table in doc.tables for row in table.rows for cell in row.cells])
        elif suffix == ".pptx":
            from pptx import Presentation
            from services.knowledge.ocr_vlm.ollama_vlm_adapter import OllamaVlmAdapter
            prs=Presentation(path); blocks=[]; adapter=OllamaVlmAdapter()
            for i, slide in enumerate(prs.slides,1):
                slide_text=[]
                for shape in slide.shapes:
                    if getattr(shape,"has_text_frame",False): slide_text.append(shape.text)
                    if getattr(shape,"shape_type",None)==13:
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                            img_path=Path(tmp.name); tmp.write(shape.image.blob)
                        try:
                            if adapter.health_check():
                                vr=adapter.invoke(str(img_path), prompt="Analyze the visible slide image. Extract text, labels, diagrams, measurements and other observable information without inventing facts.")
                                if vr.get("content"): slide_text.append(vr["content"])
                        finally: img_path.unlink(missing_ok=True)
                blocks.append(f"[Slide {i}]\n"+"\n".join(x for x in slide_text if x))
            text="\n\n".join(blocks)
        elif suffix in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook
            wb=load_workbook(path, read_only=True, data_only=True); blocks=[]
            for ws in wb.worksheets:
                rows=["\t".join("" if v is None else str(v) for v in row) for row in ws.iter_rows(values_only=True)]
                blocks.append(f"[Sheet {ws.title}]\n"+"\n".join(rows))
            text="\n\n".join(blocks)
        elif suffix in {".csv", ".tsv", ".txt", ".md", ".json", ".log", ".xml", ".html"}:
            text=p.read_text(encoding="utf-8", errors="replace")
        elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff"}:
            from services.knowledge.ocr_vlm.ollama_vlm_adapter import OllamaVlmAdapter
            if not OllamaVlmAdapter().health_check():
                raise RuntimeError("Gemma VLM is unavailable for image ingestion")
            result=OllamaVlmAdapter().invoke(path, prompt="Analyze this image for visible text, labels, measurements, findings and other observable information. Perform OCR and do not invent facts.")
            text=str(result.get("content", ""))
        else:
            raise ValueError(f"Unsupported knowledge ingestion type: {suffix or 'unknown'}")
        if not text.strip():
            return 0
        return self.ingest_text(text, source=str(path), metadata=metadata)

    def retrieve(self, query: str, top_k: int = 3, metadata_filter: dict | None = None) -> list[dict]:
        hits=self._store.search(query, top_k=top_k, metadata_filter=metadata_filter)
        return [{"claim":h["text"],"source":h["source"],"page_or_region":f"chunk_{h['chunk_index']}","confidence":float(h["score"]),"validation_state":"unverified","qdrant_point_id":h.get("point_id") } for h in hits]


def build_production_retriever() -> Retriever:
    host=os.environ.get("QDRANT_HOST","localhost")
    port=int(os.environ.get("QDRANT_PORT","6333"))
    from qdrant_client import QdrantClient
    return Retriever(VectorStore(QdrantClient(url=f"http://{host}:{port}")))


def build_in_memory_retriever() -> Retriever:
    return Retriever(MemoryVectorStore())
