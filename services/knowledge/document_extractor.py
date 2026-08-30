"""Common local extraction for PRAMAAN document inputs."""
from __future__ import annotations

import csv
import json
from pathlib import Path


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
TEXT_EXTENSIONS = {".txt", ".md", ".rst", ".log", ".xml", ".html", ".htm"}
SPREADSHEET_EXTENSIONS = {".xlsx", ".xlsm", ".xltx", ".xltm"}
PRESENTATION_EXTENSIONS = {".pptx"}
DOCUMENT_EXTENSIONS = {".docx"}


def extract_text(path_value: str | Path) -> dict:
    path = Path(path_value)
    suffix = path.suffix.lower()
    if not path.is_file():
        raise FileNotFoundError(f"Source file not found: {path}")

    if suffix in IMAGE_EXTENSIONS:
        return {"content": "", "path": str(path), "format": suffix.lstrip("."), "needs_ocr": True}

    if suffix == ".pdf":
        import fitz
        doc = fitz.open(str(path))
        pages = []
        for idx, page in enumerate(doc, start=1):
            text = (page.get_text("text") or "").strip()
            if text:
                pages.append(f"[Page {idx}]\n{text}")
        content = "\n\n".join(pages)
        return {"content": content, "path": str(path), "format": "pdf", "needs_ocr": not bool(content.strip())}

    if suffix in DOCUMENT_EXTENSIONS:
        from docx import Document
        doc = Document(str(path))
        blocks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                if any(values):
                    blocks.append(" | ".join(values))
        return {"content": "\n".join(blocks), "path": str(path), "format": "docx", "needs_ocr": False}

    if suffix in PRESENTATION_EXTENSIONS:
        from pptx import Presentation
        prs = Presentation(str(path))
        blocks = []
        for slide_no, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and getattr(shape, "text", "").strip():
                    texts.append(shape.text.strip())
            if texts:
                blocks.append(f"[Slide {slide_no}]\n" + "\n".join(texts))
        return {"content": "\n\n".join(blocks), "path": str(path), "format": "pptx", "needs_ocr": False}

    if suffix in SPREADSHEET_EXTENSIONS:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        blocks = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    rows.append(" | ".join(values))
            if rows:
                blocks.append(f"[Sheet {ws.title}]\n" + "\n".join(rows))
        return {"content": "\n\n".join(blocks), "path": str(path), "format": "xlsx", "needs_ocr": False}

    if suffix == ".csv":
        with path.open("r", encoding="utf-8", newline="") as handle:
            rows = [" | ".join(row) for row in csv.reader(handle)]
        return {"content": "\n".join(rows), "path": str(path), "format": "csv", "needs_ocr": False}

    if suffix == ".json":
        raw = path.read_text(encoding="utf-8")
        try:
            content = json.dumps(json.loads(raw), indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            content = raw
        return {"content": content, "path": str(path), "format": "json", "needs_ocr": False}

    if suffix in TEXT_EXTENSIONS:
        return {"content": path.read_text(encoding="utf-8"), "path": str(path), "format": suffix.lstrip(".") or "text", "needs_ocr": False}

    try:
        return {"content": path.read_text(encoding="utf-8"), "path": str(path), "format": suffix.lstrip(".") or "text", "needs_ocr": False}
    except UnicodeDecodeError as exc:
        raise ValueError(
            f"Unsupported binary input '{suffix}'. PRAMAAN supports PDF, DOCX, PPTX, XLSX, CSV, JSON, TXT/Markdown, and common image formats."
        ) from exc
