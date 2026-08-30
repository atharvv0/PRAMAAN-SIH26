from __future__ import annotations

import json
from pathlib import Path
from typing import Any


def extract_document(path: str) -> dict[str, Any]:
    """Extract searchable text from common PRAMAAN document formats.

    Returns {content, source, pages?}. Images are intentionally handled by the
    VLM/OCR tool instead of being converted to meaningless raw bytes here.
    """
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        import fitz
        doc = fitz.open(path)
        pages = [page.get_text("text") for page in doc]
        return {"content": "\n\n".join(f"[Page {i+1}]\n{t}" for i, t in enumerate(pages) if t.strip()), "pages": pages}

    if suffix in {".docx", ".doc"}:
        if suffix == ".doc":
            raise ValueError("Legacy .doc is not supported; convert it to .docx first.")
        from docx import Document
        doc = Document(path)
        parts = [para.text for para in doc.paragraphs if para.text.strip()]
        for table_index, table in enumerate(doc.tables, 1):
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            parts.append(f"[Table {table_index}]\n" + "\n".join(rows))
        return {"content": "\n\n".join(parts)}

    if suffix in {".pptx", ".ppt"}:
        if suffix == ".ppt":
            raise ValueError("Legacy .ppt is not supported; convert it to .pptx first.")
        from pptx import Presentation
        prs = Presentation(path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            bits: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        bits.append(text)
            slides.append(f"[Slide {i}]\n" + "\n".join(bits))
        return {"content": "\n\n".join(slides), "pages": slides}

    if suffix in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                if any(values):
                    rows.append(" | ".join(values))
            if rows:
                parts.append(f"[Sheet {ws.title}]\n" + "\n".join(rows))
        return {"content": "\n\n".join(parts)}

    if suffix == ".csv":
        return {"content": p.read_text(encoding="utf-8", errors="replace")}

    if suffix == ".json":
        return {"content": json.dumps(json.loads(p.read_text(encoding="utf-8")), indent=2, ensure_ascii=False)}

    if suffix in {".txt", ".md", ".markdown", ".log", ".xml", ".html", ".htm"}:
        return {"content": p.read_text(encoding="utf-8", errors="replace")}

    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".tif"}:
        return {"content": "", "requires_vision": True}

    raise ValueError(f"Unsupported document type: {suffix or 'unknown'}")
