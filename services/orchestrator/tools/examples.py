"""
DEMO-ONLY tools used to prove the orchestration loop end-to-end (master prompt
section 30, "First Demo For My Module"). These are NOT the real production tools.

Real document/OCR/RAG tools belong to services/knowledge. Real sandboxed code
execution belongs to services/sandbox. These two exist only so the Planner ->
Executor -> ToolRegistry loop can be demonstrated and tested before those modules
land. Do not extend this file with real capabilities — register real tools in their
owning service and add them to the shared registry in tools/registry_instance.py.
"""
from __future__ import annotations

from pathlib import Path

from services.orchestrator.tools.base import ToolAdapter


class ReadFileTool(ToolAdapter):
    """Reads a local text file. id: file.read (matches docs/agent-contract.md
    tool id examples). No sandboxing, no permission check yet — that's
    services/governance/policy_engine, not wired in until Phase 7."""

    id = "file.read"
    required_permissions = ["file.read"]
    declares_network_access = False

    def invoke(self, inputs: dict) -> dict:
        path = inputs.get("path")
        if not path:
            raise ValueError("ReadFileTool requires 'path' in inputs")
        p = Path(path)
        suffix = p.suffix.lower()
        text = ""
        if suffix == ".pdf":
            import fitz
            doc = fitz.open(path)
            text = "\n\n".join(f"[Page {i+1}]\n{page.get_text('text')}" for i, page in enumerate(doc))
        elif suffix == ".docx":
            from docx import Document
            doc = Document(path)
            text = "\n".join([para.text for para in doc.paragraphs] + [cell.text for table in doc.tables for row in table.rows for cell in row.cells])
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            slides=[]
            for i, slide in enumerate(prs.slides, 1):
                lines=[]
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        lines.append(shape.text)
                slides.append(f"[Slide {i}]\n" + "\n".join(x for x in lines if x))
            text="\n\n".join(slides)
        elif suffix in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook
            wb=load_workbook(path, read_only=True, data_only=True)
            blocks=[]
            for ws in wb.worksheets:
                rows=["\t".join("" if v is None else str(v) for v in row) for row in ws.iter_rows(values_only=True)]
                blocks.append(f"[Sheet {ws.title}]\n"+"\n".join(rows))
            text="\n\n".join(blocks)
        elif suffix in {".csv", ".tsv"}:
            text=p.read_text(encoding="utf-8", errors="replace")
        elif suffix in {".json"}:
            text=p.read_text(encoding="utf-8", errors="replace")
        elif suffix in {".txt", ".md", ".log", ".xml", ".html"}:
            text=p.read_text(encoding="utf-8", errors="replace")
        else:
            raise ValueError(f"Unsupported text-readable file type: {suffix or 'unknown'}")
        return {
            "content": text,
            "path": path,
            "evidence": [{
                "claim": f"Source document content extracted successfully from {p.name}.",
                "source": path,
                "page_or_region": "page_1",
                "confidence": 1.0,
                "validation_state": "unverified",
            }],
        }


class OcrProcessNaiveTool(ToolAdapter):
    """DEMO-ONLY placeholder for the real OCR/VLM pipeline. id:
    ocr.process_naive (matches the "ocr.process" tool id family in
    docs/agent-contract.md / architecture.md). It does not actually run OCR —
    it reads a plain text file and fabricates an EvidenceRecord-shaped citation
    so the executor's evidence-population path (docs/agent-contract.md
    "EvidenceRecord") can be proven end-to-end before the real multimodal
    pipeline exists. TODO: replace with the real services/knowledge OCR/VLM
    pipeline — keep the output shape (content + evidence[]) stable so the
    plans referencing this tool id don't need to change, just the
    registration in tools/registry_instance.py."""

    id = "ocr.process_naive"
    required_permissions = ["file.read"]
    declares_network_access = False

    def invoke(self, inputs: dict) -> dict:
        path = inputs.get("path")
        if not path:
            raise ValueError("OcrProcessNaiveTool requires 'path' in inputs")
        text = Path(path).read_text(encoding="utf-8")
        return {
            "content": text,
            "path": path,
            "evidence": [
                {
                    "claim": "Text extracted from source document (placeholder OCR — not real vision/OCR)",
                    "source": path,
                    "page_or_region": "page_1",
                    "confidence": 0.5,
                    "validation_state": "unverified",
                }
            ],
        }


class SummarizeTextTool(ToolAdapter):
    """Naive placeholder summarizer: first 3 sentences, no model call. id:
    text.summarize_naive. TODO(Phase 5): replace with a real call through
    services/model_control once the Model Router exists — keep this tool's id
    stable or update every plan that references it."""

    id = "text.summarize_naive"
    required_permissions = []
    declares_network_access = False

    def invoke(self, inputs: dict) -> dict:
        content = None
        for value in inputs.values():
            if isinstance(value, dict) and "content" in value:
                content = value["content"]
                break
        if content is None:
            raise ValueError("SummarizeTextTool found no upstream 'content' to summarize")

        sentences = [s.strip() for s in content.replace("\n", " ").split(".") if s.strip()]
        summary = ". ".join(sentences[:3])
        if summary and not summary.endswith("."):
            summary += "."
        return {"summary": summary, "sentence_count": len(sentences)}


class NetworkFetchDemoTool(ToolAdapter):
    """DEMO-ONLY tool that declares network access, used purely to prove the
    sovereignty story live: the DefaultPolicyEngine (services/governance) denies
    any tool with declares_network_access=True before it ever runs — this tool's
    invoke() is intentionally never reached in a correctly-configured sovereign
    deployment. id: network.fetch_demo."""

    id = "network.fetch_demo"
    required_permissions = ["network.egress"]
    declares_network_access = True

    def invoke(self, inputs: dict) -> dict:
        # Never actually reached under DefaultPolicyEngine — if this line runs,
        # the sovereignty boundary has a hole in it.
        raise RuntimeError("network.fetch_demo should never execute in sovereign mode")
