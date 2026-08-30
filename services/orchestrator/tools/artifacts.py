from __future__ import annotations

import csv
import json
import mimetypes
import os
from pathlib import Path
from uuid import uuid4

from services.orchestrator.tools.base import ToolAdapter


ARTIFACT_ROOT = Path(__file__).resolve().parents[3] / "data" / "deliverables"
ARTIFACT_ROOT.mkdir(parents=True, exist_ok=True)


def _clean_thinking(text: str) -> str:
    cleaned = str(text or "")
    cleaned = re.sub(r"<think>.*?</think>", "", cleaned, flags=re.I | re.S)
    cleaned = re.sub(r"<analysis>.*?</analysis>", "", cleaned, flags=re.I | re.S)
    return cleaned.strip()


def _content_from_inputs(inputs: dict) -> str:
    for key in ("content", "response", "answer", "summary"):
        value = inputs.get(key)
        if isinstance(value, str) and value.strip():
            return _clean_thinking(value)
    for key, value in inputs.items():
        if not str(key).startswith("upstream_") or not isinstance(value, dict):
            continue
        for field in ("content", "response", "answer", "summary"):
            candidate = value.get(field)
            if isinstance(candidate, str) and candidate.strip():
                return _clean_thinking(candidate)
    raise ValueError("artifact.write requires upstream content or a direct content field")


def _normalise_format(value: str | None, filename: str | None) -> str:
    raw = (value or "").strip().lower().lstrip(".")
    if raw:
        return {"text": "txt", "markdown": "md", "word": "docx", "powerpoint": "pptx", "excel": "xlsx"}.get(raw, raw)
    suffix = Path(filename or "").suffix.lower().lstrip(".")
    return suffix or "txt"


class ArtifactWriteTool(ToolAdapter):
    id = "artifact.write"
    required_permissions = ["file.write"]
    declares_network_access = False

    def invoke(self, inputs: dict) -> dict:
        content = _content_from_inputs(inputs)
        intent = str(inputs.get("intent") or "")
        import re
        match = re.search(r"\b(\d+)\s*(?:-| )?lines?\b", intent, flags=re.I)
        fmt = _normalise_format(inputs.get("format"), inputs.get("filename"))
        if match:
            required_lines = int(match.group(1))
            lines = content.splitlines()
            if len(lines) != required_lines:
                raise ValueError(f"Requested exactly {required_lines} lines, but generated output contains {len(lines)} lines")
        requested_name = str(inputs.get("filename") or f"pramaan_output.{fmt}")
        safe_name = Path(requested_name).name or f"pramaan_output.{fmt}"
        if not Path(safe_name).suffix:
            safe_name = f"{safe_name}.{fmt}"
        path = ARTIFACT_ROOT / f"{uuid4()}_{safe_name}"

        if fmt in {"txt", "md", "json", "csv"}:
            if fmt == "json":
                try:
                    parsed = json.loads(content)
                    path.write_text(json.dumps(parsed, indent=2, ensure_ascii=False), encoding="utf-8")
                except json.JSONDecodeError:
                    path.write_text(content, encoding="utf-8")
            elif fmt == "csv":
                path.write_text(content, encoding="utf-8")
            else:
                path.write_text(content, encoding="utf-8")
        elif fmt == "docx":
            from docx import Document
            doc = Document()
            for line in content.splitlines() or [""]:
                doc.add_paragraph(line)
            doc.save(path)
        elif fmt == "pdf":
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
            c = canvas.Canvas(str(path), pagesize=A4)
            width, height = A4
            y = height - 48
            for raw_line in content.splitlines() or [""]:
                line = raw_line
                while len(line) > 110:
                    c.drawString(40, y, line[:110])
                    y -= 14
                    line = line[110:]
                    if y < 48:
                        c.showPage(); y = height - 48
                c.drawString(40, y, line)
                y -= 14
                if y < 48:
                    c.showPage(); y = height - 48
            c.save()
        elif fmt == "pptx":
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            for idx, line in enumerate(content.splitlines() or [""]):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"PRAMAAN Output {idx + 1}"
                slide.placeholders[1].text = line
            if len(prs.slides) == 0:
                prs.slides.add_slide(prs.slide_layouts[1])
            prs.save(path)
        elif fmt == "xlsx":
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "PRAMAAN Output"
            for row_idx, line in enumerate(content.splitlines() or [""], 1):
                ws.cell(row=row_idx, column=1, value=line)
            wb.save(path)
        else:
            raise ValueError(f"Unsupported artifact format: {fmt}")

        mime = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        return {
            "artifact": True,
            "path": str(path),
            "filename": safe_name,
            "format": fmt,
            "mime_type": mime,
            "size_bytes": path.stat().st_size,
        }
