from __future__ import annotations

import csv
import json
import re
import textwrap
from pathlib import Path
from uuid import uuid4

from ..db.database import FileRecord, Task, session_scope
from ..db.repository import UPLOAD_ROOT, repo

ARTIFACT_ROOT = UPLOAD_ROOT.parent / "deliverables"
ARTIFACT_ROOT.mkdir(parents=True, exist_ok=True)

_FORMAT_EXTENSIONS = {
    "txt": ".txt",
    "text": ".txt",
    "docx": ".docx",
    "pdf": ".pdf",
    "pptx": ".pptx",
    "ppt": ".pptx",
    "xlsx": ".xlsx",
    "xls": ".xlsx",
    "csv": ".csv",
    "json": ".json",
}


def _clean_text(value: object) -> str:
    text = str(value or "").strip()
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r"^\s*(?:thinking|analysis)\s*:\s*.*?(?=\n\s*(?:answer|final|response)\s*:|$)", "", text, flags=re.IGNORECASE | re.DOTALL)
    text = re.sub(r"^\s*(?:answer|final answer|response)\s*:\s*", "", text, flags=re.IGNORECASE)
    return text.strip()


def requested_format(intent: str) -> str | None:
    lower = intent.lower()
    for fmt, ext in _FORMAT_EXTENSIONS.items():
        if re.search(rf"(?:\.|\b){re.escape(fmt)}\b", lower) or f"{ext} file" in lower:
            return "txt" if fmt == "text" else ("xlsx" if fmt == "xls" else fmt)
    return None


def requested_line_count(intent: str) -> int | None:
    match = re.search(r"\b(?:in|within|exactly)\s+(\d+)\s+lines?\b", intent.lower())
    return int(match.group(1)) if match else None


def requested_artifact(intent: str) -> bool:
    lower = intent.lower()
    return bool(
        requested_format(intent)
        or re.search(r"\b(?:file|deliverable|export|save|download|generate|give me)\b", lower)
    )


def _exact_line_count(text: str, count: int) -> str:
    if count <= 0:
        return text.strip()
    cleaned = _clean_text(text)
    if not cleaned:
        return cleaned

    # Prefer existing non-empty lines, then sentence boundaries, then deterministic wrapping.
    lines = [re.sub(r"^\s*[-*•]\s*", "", x).strip() for x in cleaned.splitlines() if x.strip()]
    if len(lines) < count:
        sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+|;\s+", " ".join(lines)) if s.strip()]
        lines = sentences
    if len(lines) < count:
        words = " ".join(lines).split()
        width = max(20, round(len(" ".join(lines)) / count))
        lines = textwrap.wrap(" ".join(words), width=width, break_long_words=False, break_on_hyphens=False)

    if len(lines) > count:
        groups: list[list[str]] = [[] for _ in range(count)]
        for i, line in enumerate(lines):
            groups[min(count - 1, (i * count) // len(lines))].append(line)
        lines = [" ".join(group).strip() for group in groups]
        lines = [line for line in lines if line]

    # Never invent content. If wrapping yielded fewer lines than requested, split the longest lines.
    while len(lines) < count:
        idx = max(range(len(lines)), key=lambda i: len(lines[i]))
        words = lines[idx].split()
        if len(words) < 2:
            break
        mid = max(1, len(words) // 2)
        lines[idx:idx + 1] = [" ".join(words[:mid]), " ".join(words[mid:])]

    if len(lines) < count:
        # Extremely short source: keep the truthful content and do not fabricate lines.
        return "\n".join(lines)
    return "\n".join(lines[:count])


def _write_docx(path: Path, content: str, title: str) -> None:
    from docx import Document
    doc = Document()
    doc.add_heading(title or "PRAMAAN Deliverable", level=0)
    for line in content.splitlines() or [content]:
        doc.add_paragraph(line)
    doc.save(path)


def _write_pdf(path: Path, content: str, title: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=A4)
    story = [Paragraph(title or "PRAMAAN Deliverable", styles["Title"]), Spacer(1, 12)]
    for line in content.splitlines() or [content]:
        story.append(Paragraph(line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), styles["BodyText"]))
        story.append(Spacer(1, 6))
    doc.build(story)


def _write_pptx(path: Path, content: str, title: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title or "PRAMAAN Deliverable"
    title_slide.placeholders[1].text = "Generated locally by PRAMAAN"

    lines = [x.strip() for x in content.splitlines() if x.strip()]
    chunk_size = 8
    for idx in range(0, max(1, len(lines)), chunk_size):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Summary {idx // chunk_size + 1}"
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.2))
        frame = box.text_frame
        frame.clear()
        for line in lines[idx:idx + chunk_size] or [content]:
            p = frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
    prs.save(path)


def _write_xlsx(path: Path, content: str, title: str) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "PRAMAAN"
    ws["A1"] = title or "PRAMAAN Deliverable"
    for row, line in enumerate(content.splitlines() or [content], start=3):
        ws.cell(row=row, column=1, value=line)
    wb.save(path)


def _write_csv(path: Path, content: str, title: str) -> None:
    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.writer(handle)
        writer.writerow([title or "PRAMAAN Deliverable"])
        for line in content.splitlines() or [content]:
            writer.writerow([line])


def _write_json(path: Path, content: str, title: str) -> None:
    path.write_text(json.dumps({"title": title, "content": content.splitlines() or [content]}, indent=2, ensure_ascii=False), encoding="utf-8")


def generate_task_deliverable(task_id: str, state: dict, title: str, intent: str) -> dict | None:
    """Create and persist a requested user artifact. Returns metadata or None."""
    if not requested_artifact(intent):
        return None

    final = state.get("final_output") or {}
    response = _clean_text(final.get("response") if isinstance(final, dict) else final)
    if not response:
        return None

    fmt = requested_format(intent) or "txt"
    line_count = requested_line_count(intent)

    # For summary requests, prefer the model-backed summary itself for the artifact.
    if "summar" in intent.lower():
        for output in reversed(state.get("tool_outputs", [])):
            if isinstance(output, dict) and isinstance(output.get("summary"), str) and output["summary"].strip():
                response = _clean_text(output["summary"])
                break

    if line_count:
        response = _exact_line_count(response, line_count)

    artifact_id = str(uuid4())
    ext = _FORMAT_EXTENSIONS.get(fmt, ".txt")
    path = ARTIFACT_ROOT / f"{artifact_id}_{Path(title or 'PRAMAAN_Deliverable').stem}{ext}"

    if fmt == "docx":
        _write_docx(path, response, title)
    elif fmt == "pdf":
        _write_pdf(path, response, title)
    elif fmt == "pptx":
        _write_pptx(path, response, title)
    elif fmt == "xlsx":
        _write_xlsx(path, response, title)
    elif fmt == "csv":
        _write_csv(path, response, title)
    elif fmt == "json":
        _write_json(path, response, title)
    else:
        path.write_text(response.rstrip() + "\n", encoding="utf-8")
        fmt = "txt"

    data = path.read_bytes()
    with session_scope() as s:
        task = s.get(Task, task_id)
        if task is None:
            raise ValueError(f"Task not found: {task_id}")
        rec = FileRecord(
            file_id=str(uuid4()),
            project_id=task.project_id,
            uploaded_by=task.created_by,
            filename=path.name,
            mime_type={
                "txt": "text/plain",
                "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "pdf": "application/pdf",
                "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "csv": "text/csv",
                "json": "application/json",
            }.get(fmt, "application/octet-stream"),
            size_bytes=len(data),
            storage_path=str(path),
            sha256=__import__("hashlib").sha256(data).hexdigest(),
            sensitivity_class=task.sensitivity_class,
        )
        s.add(rec)
        s.commit()
        s.refresh(rec)
        deliverable_id = repo.create_deliverable(task_id, rec.file_id, fmt, "approved")

    return {
        "id": deliverable_id,
        "file_id": rec.file_id,
        "name": rec.filename,
        "type": fmt,
        "path": str(path),
        "downloadUrl": f"/api/v1/files/{rec.file_id}/download",
        "lineCount": len([x for x in response.splitlines() if x.strip()]),
        "source": "final_output",
    }
